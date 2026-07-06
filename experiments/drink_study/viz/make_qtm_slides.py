"""Generate methodology + results figures (and a starter .pptx) for the
QTM-mocap vs video cup-tracking accuracy comparison.

Story (slides): how accurately the video cup track matches sub-mm QTM mocap GT
across 669 drinking reps. The duration->geometry re-pairing fix is a footnote.

    python make_qtm_slides.py            # -> slides/*.png + slides/qtm_accuracy.pptx

All figures read cache/qtm_align.json (per-rep inlier_rms_mm, frac_fail, cls).
The two methodology trajectory panels recompute one clean + one localized rep
from the C3D + video track via qtm_align (needs ezc3d + the local caches).
"""
from __future__ import annotations
import sys as _s, pathlib as _p  # drink_study lib path shim
for _q in _p.Path(__file__).resolve().parents:
    if (_q / 'lib' / 'segment_cup_only.py').exists():
        _s.path.insert(0, str(_q / 'lib')); _s.path.insert(0, str(_q)); _s.path.insert(0, str(_q.parents[1])); break
import json
from pathlib import Path
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Patch

import qtm_align as Q
from qtm_c3d import load_trial
import segment_cup_only as S

HERE = Path(__file__).resolve().parent
from _paths import CACHE
OUT = HERE / "slides"
OUT.mkdir(exist_ok=True)

# palette
C_CLEAN, C_LOCAL, C_BROKEN = "#2a9d8f", "#e9c46a", "#e76f51"
C_MOCAP, C_VIDEO = "#264653", "#e76f51"
plt.rcParams.update({"font.size": 12, "axes.spLines.top" if False else "axes.spines.top": False,
                     "axes.spines.right": False, "figure.dpi": 140})

d = json.load(open(CACHE / "qtm_align.json"))
reps = [r for p in d.values() if isinstance(p, dict) and p.get("ok") for r in p["reps"]]
ir = np.array([r["inlier_rms_mm"] for r in reps])
ff = np.array([r["frac_fail"] for r in reps])
cls = np.array([r["cls"] for r in reps])
N = len(reps)
n_clean = int((cls == "clean").sum())
n_local = int((cls == "localized").sum())
n_broken = int((cls == "broken").sum())
print(f"{N} reps: {n_clean} clean / {n_local} localized / {n_broken} broken")


def save(fig, name):
    p = OUT / name
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("wrote", p)
    return p


# ---------------------------------------------------------------- methodology
def fig_pipeline():
    """Block diagram of the comparison pipeline (vertical inputs -> shared chain)."""
    fig, ax = plt.subplots(figsize=(12, 3.4))
    ax.axis("off")
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    boxes = [
        ("QTM mocap\n4 cup markers\nsub-mm GT", C_MOCAP, "white"),
        ("Video cup track\n10-cam consensus\nKF / RTS", C_VIDEO, "white"),
        ("Pair rep ↔ take\norder-preserving\nKabsch-RMS", "#8ab17d", "black"),
        ("Temporal sync\nspeed-curve\nxcorr", "#8ab17d", "black"),
        ("Robust Kabsch fit\ninlier RMS\n+ frac_fail", "#8ab17d", "black"),
        ("Classify rep\nclean / localized\n/ broken", "#577590", "white"),
    ]
    n = len(boxes)
    w, gap = 0.142, 0.025
    x0 = (1 - (n * w + (n - 1) * gap)) / 2
    x = x0
    centers = []
    for i, (t, fc, tc) in enumerate(boxes):
        ax.add_patch(plt.Rectangle((x, 0.32), w, 0.46, fc=fc, ec="none"))
        ax.text(x + w / 2, 0.55, t, ha="center", va="center", color=tc,
                fontsize=9.5, weight="bold", linespacing=1.4)
        centers.append(x + w / 2)
        if i < n - 1:
            ax.annotate("", xy=(x + w + gap, 0.55), xytext=(x + w, 0.55),
                        arrowprops=dict(arrowstyle="-|>", color="#333", lw=2.2))
        x += w + gap
    ax.text(0.5, 0.10, "Each video drinking rep is matched to its mocap take, time-aligned, "
            "then rigidly fit; agreement = tracking accuracy.",
            ha="center", fontsize=10.5, style="italic", color="#555")
    fig.suptitle("Methodology: QTM mocap as ground truth for video cup tracking",
                 fontsize=15, weight="bold", y=1.0)
    return save(fig, "fig1_pipeline.png")


def _track_vs_mocap(video_name, c3d_stem):
    vpath = CACHE / "track3d_clean3d_refill" / video_name
    r = Q._synced_rep(vpath, c3d_stem)
    mo, v, lag, corr = r
    R, t, rms = Q.kabsch(mo, v, robust=True)
    moT = mo @ R.T + t
    err = np.linalg.norm(v - moT, axis=1)
    return moT, v, err, rms


def fig_method_example():
    """Two panels: a clean rep and a localized rep — per-frame error with the
    cup-only phase segmentation shaded behind it (failure lands in 'drinking')."""
    clean = ("P02_P02_drinking_left_20240110_111516__clean3d_refill.json", "P02_0036")
    local = ("P02_P02_drinking_left_20240110_111430__clean3d_refill.json", "P02_0033")
    fig, axes = plt.subplots(1, 2, figsize=(11.5, 4.0), sharey=True)
    for ax, (vn, cs), title in zip(axes, (clean, local), ("Clean rep", "Localized rep (drink apex fails)")):
        moT, v, err, rms = _track_vs_mocap(vn, cs)
        # phases from the OMC (mocap) trajectory, on the same COMMON_HZ grid as err
        tr = load_trial(cs)
        mr = Q._resample(tr.centroid(), tr.rate)
        seg = S.segment_cup_only(mr, fps=Q.COMMON_HZ)
        tt = np.arange(len(err)) / Q.COMMON_HZ
        ff_ = float(np.mean(err > Q.FAIL_MM))
        T = len(err)
        for nm, s, e in seg["intervals"]:
            ax.axvspan(min(s, T) / Q.COMMON_HZ, min(e, T) / Q.COMMON_HZ,
                       color=S.PHASE_COLORS[nm], alpha=0.22, lw=0)
        ax.plot(tt, err, color="#222", lw=1.8)
        ax.axhline(Q.FAIL_MM, color="#666", ls="--", lw=1)
        ax.set_title(f"{title}\ninlier RMS {rms:.1f} mm   frac_fail {ff_:.2f}", fontsize=11)
        ax.set_xlabel("time (s)")
        ax.set_ylim(0, max(80, err.max() * 1.1))
        ax.set_xlim(0, tt[-1])
    axes[0].set_ylabel("mocap–video distance (mm)")
    handles = [Patch(color=S.PHASE_COLORS[p], alpha=0.5, label=p) for p in S.PHASE_NAMES]
    handles.append(plt.Line2D([0], [0], color="#666", ls="--", label="50 mm fail"))
    fig.legend(handles=handles, loc="upper center", ncol=6, fontsize=8.5, frameon=False,
               bbox_to_anchor=(0.5, 1.02))
    fig.suptitle("Per-frame error vs drink phase — failure is confined to the drinking dwell",
                 fontsize=13, weight="bold", y=0.93)
    fig.tight_layout(rect=[0, 0, 1, 0.88])
    return save(fig, "fig2_method_example.png")


def fig_fracfail_by_phase():
    pf = json.load(open(CACHE / "phase_failure.json"))
    phases = pf["phases"]
    vals = [pf["frac_fail_by_phase"][p] for p in phases]
    fig, ax = plt.subplots(figsize=(8.5, 4.5))
    bars = ax.bar(range(len(phases)), vals,
                  color=[S.PHASE_COLORS[p] for p in phases], edgecolor="#333")
    for b, p in zip(bars, phases):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.004,
                f"{pf['frac_fail_by_phase'][p]:.1%}\n{pf['median_err_by_phase'][p]:.0f} mm",
                ha="center", fontsize=9.5)
    ax.set_xticks(range(len(phases)))
    ax.set_xticklabels([p.replace("_", "\n") for p in phases], fontsize=10)
    ax.set_ylabel("frac_fail  (fraction of frames >50 mm off)")
    ax.set_ylim(0, max(vals) * 1.25)
    ax.set_title(f"Tracking failure is concentrated in the drinking phase  "
                 f"(N={pf['n_reps']} reps)", weight="bold")
    return save(fig, "fig7_fracfail_by_phase.png")


def fig_within_phase():
    """One continuous failure-rate line across the whole drinking movement
    (forward_transport -> drinking -> back_transport mapped onto [0,1]), with the
    phase regions shaded so it reads as a single occlusion arc."""
    pf = json.load(open(CACHE / "phase_failure.json"))
    nb = pf["movement_nbins"]
    x = (np.arange(nb) + 0.5) / nb
    prof = np.array(pf["movement_fail_rate"])
    b1 = pf["movement_bound_fwd_drink"]
    b2 = pf["movement_bound_drink_back"]
    fig, ax = plt.subplots(figsize=(9.5, 4.6))
    # phase background bands
    ax.axvspan(0, b1, color=S.PHASE_COLORS["forward_transport"], alpha=0.18, lw=0)
    ax.axvspan(b1, b2, color=S.PHASE_COLORS["drinking"], alpha=0.18, lw=0)
    ax.axvspan(b2, 1, color=S.PHASE_COLORS["back_transport"], alpha=0.18, lw=0)
    ax.plot(x, prof, color="#222", lw=2.6)
    ax.fill_between(x, 0, prof, color="#444", alpha=0.08)
    ymax = prof.max() * 1.18
    for b, lab in ((b1, "grasp→drink"), (b2, "drink→release")):
        ax.axvline(b, color="#666", ls=":", lw=1.2)
    ax.text(b1 / 2, ymax * 0.95, "forward\ntransport", ha="center", va="top", fontsize=10,
            color="#2f6ea8")
    ax.text((b1 + b2) / 2, ymax * 0.95, "drinking\n(apex)", ha="center", va="top",
            fontsize=10, color="#b03a2e", weight="bold")
    ax.text((b2 + 1) / 2, ymax * 0.95, "back\ntransport", ha="center", va="top", fontsize=10,
            color="#b9770e")
    ax.set_xlabel("normalised position through the drinking movement  (0 = grasp → 1 = release)")
    ax.set_ylabel("failure rate  (fraction of frames >50 mm)")
    ax.set_title("Failure is one continuous arc — rises to the mouth, peaks at the apex, decays",
                 weight="bold", fontsize=12.5)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, ymax)
    return save(fig, "fig8_within_phase.png")


# -------------------------------------------------------------------- results
def fig_rms_hist():
    fig, ax = plt.subplots(figsize=(8, 4.2))
    ax.hist(ir, bins=np.arange(0, ir.max() + 1, 0.5), color=C_MOCAP, alpha=.85)
    med = np.median(ir)
    ax.axvline(med, color=C_VIDEO, lw=2, label=f"median {med:.1f} mm")
    ax.axvline(np.percentile(ir, 90), color="#888", ls="--", lw=1.5,
               label=f"p90 {np.percentile(ir,90):.1f} mm")
    ax.set_xlabel("per-rep inlier RMS (mm)  —  tracking fidelity vs mocap")
    ax.set_ylabel("number of reps")
    ax.set_title(f"Video cup track agrees with mocap to a few mm  (N={N} reps)",
                 weight="bold")
    ax.legend()
    return save(fig, "fig3_rms_hist.png")


def fig_class_bar():
    fig, ax = plt.subplots(figsize=(6.5, 4.2))
    vals = [n_clean, n_local, n_broken]
    labs = [f"clean\n{n_clean} ({n_clean/N:.0%})",
            f"localized\n{n_local} ({n_local/N:.0%})",
            f"broken\n{n_broken} ({n_broken/N:.0%})"]
    bars = ax.bar(labs, vals, color=[C_CLEAN, C_LOCAL, C_BROKEN])
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 4, str(v), ha="center", weight="bold")
    ax.set_ylabel("number of reps")
    ax.set_title("Rep accuracy classes  —  0 broken after geometry re-pairing", weight="bold")
    ax.set_ylim(0, max(vals) * 1.15)
    return save(fig, "fig4_class_bar.png")


def fig_scatter():
    fig, ax = plt.subplots(figsize=(7.5, 5))
    colors = {"clean": C_CLEAN, "localized": C_LOCAL, "broken": C_BROKEN}
    for c in ("clean", "localized", "broken"):
        m = cls == c
        ax.scatter(ir[m], ff[m], s=22, alpha=.6, color=colors[c], label=c, edgecolor="none")
    ax.axhline(0.10, color="#999", ls="--", lw=1)
    ax.axvline(15, color="#999", ls="--", lw=1)
    ax.text(15.2, 0.35, "15 mm fidelity\nthreshold", fontsize=8, color="#666")
    ax.text(0.3, 0.105, "10% frac_fail threshold", fontsize=8, color="#666")
    ax.set_xlabel("inlier RMS (mm)  —  fidelity")
    ax.set_ylabel("frac_fail  —  fraction of trajectory >50 mm off")
    ax.set_title("Every rep tracks tightly (<15 mm); class = how much of it fails", weight="bold")
    ax.legend()
    return save(fig, "fig5_scatter.png")


def fig_per_participant():
    rows = [(p["participant"],
             np.median([r["inlier_rms_mm"] for r in p["reps"]]),
             len(p["reps"]))
            for p in d.values() if isinstance(p, dict) and p.get("ok")]
    rows.sort(key=lambda x: x[1])
    pids = [r[0] for r in rows]
    meds = [r[1] for r in rows]
    fig, ax = plt.subplots(figsize=(11, 4.2))
    ax.bar(pids, meds, color=C_MOCAP)
    ax.axhline(np.median(ir), color=C_VIDEO, lw=1.5, ls="--",
               label=f"cohort median {np.median(ir):.1f} mm")
    ax.set_ylabel("median inlier RMS (mm)")
    ax.set_xlabel("participant (sorted)")
    ax.set_title("Tracking accuracy is consistent across all 22 participants", weight="bold")
    ax.tick_params(axis="x", rotation=60, labelsize=8)
    ax.legend()
    return save(fig, "fig6_per_participant.png")


def fig_failure_factors():
    """Untangled drivers of failure, three panels:
      (a) univariate vs adjusted odds — how much each effect shrinks when the
          correlated predictors are controlled for;
      (b) hand-on-cup as a GATE — failure only when the hand is on the cup;
      (c) speed is NON-MONOTONIC — slow fails at the mouth, fast in transport."""
    ff = json.load(open(CACHE / "failure_factors.json"))
    keys = ff["factors"]
    short = {keys[0]: "cup near\nmouth", keys[1]: "hand on\ncup", keys[2]: "speed"}
    fig, (a, b, c) = plt.subplots(1, 3, figsize=(13.5, 4.4),
                                  gridspec_kw=dict(width_ratios=[1.2, 0.8, 1.1]))
    cols = [C_VIDEO, "#e9a23b", "#577590"]

    # (a) univariate vs adjusted grouped bars
    labels = [short[k] for k in keys]
    uni = [ff["univariate_odds"][k] for k in keys]
    adj = [ff["adjusted_odds"][k] for k in keys]
    ci = ff.get("odds_ci95", {})
    # asymmetric error bars (adjusted odds -> 95% CI) for each factor
    yerr = np.array([[adj[i] - ci[keys[i]][0] for i in range(3)],
                     [ci[keys[i]][1] - adj[i] for i in range(3)]]) if ci else None
    x = np.arange(3); w = 0.38
    a.bar(x - w / 2, uni, w, color="#bbb", edgecolor="#333", label="alone (univariate)")
    a.bar(x + w / 2, adj, w, color=cols, edgecolor="#333", label="adjusted (all 3)")
    if yerr is not None:
        a.errorbar(x + w / 2, adj, yerr=yerr, fmt="none", ecolor="#222", capsize=4, lw=1.5,
                   label="95% CI (cluster boot)")
    for xi, (u, ad) in enumerate(zip(uni, adj)):
        a.text(xi - w / 2, u + 0.15, f"×{u:.1f}", ha="center", fontsize=9, color="#555")
    a.axhline(1, color="#999", ls="--", lw=1)
    a.set_xticks(x); a.set_xticklabels(labels, fontsize=10)
    a.set_ylabel("failure odds per +1 SD")
    top = max(max(uni), max([ci[keys[i]][1] for i in range(3)]) if ci else 0)
    a.set_ylim(0, top * 1.12)
    a.set_title("Effects shrink once correlated\npredictors are controlled", fontsize=11, weight="bold")
    a.legend(fontsize=8, loc="upper right")

    # (b) hand-on gate
    g = ff["gate"]
    bars = b.bar(["hand\nON cup", "hand\nOFF cup"],
                 [g["fail_on"], g["fail_off"]], color=["#e9a23b", "#cccccc"], edgecolor="#333")
    for bar, v in zip(bars, [g["fail_on"], g["fail_off"]]):
        b.text(bar.get_x() + bar.get_width() / 2, v + 0.005, f"{v:.0%}", ha="center", weight="bold")
    b.set_ylabel("failure rate")
    b.set_ylim(0, max(g["fail_on"], 0.05) * 1.25)
    b.set_title("Hand-on-cup is a GATE\n(fail only when hand on)", fontsize=11, weight="bold")

    # (c) speed non-monotonic interaction
    inter = ff["speed_mouth_interaction"]
    xb = np.arange(len(inter)); w2 = 0.38
    slow = [r["fail_slow"] for r in inter]; fast = [r["fail_fast"] for r in inter]
    c.bar(xb - w2 / 2, slow, w2, color="#4c78c8", edgecolor="#333", label="slow")
    c.bar(xb + w2 / 2, fast, w2, color="#e45756", edgecolor="#333", label="fast")
    c.set_xticks(xb)
    c.set_xticklabels([f"{r['band']}\nmm" for r in inter], fontsize=9)
    c.set_xlabel("cup → mouth distance")
    c.set_ylabel("failure rate")
    c.set_title("Speed effect flips with distance:\nslow@mouth, fast@transport", fontsize=11, weight="bold")
    c.legend(title="cup speed", fontsize=9)

    fig.suptitle(f"Untangling the drivers of failure  (logistic AUC {ff['auc']:.2f}, "
                 f"N={ff['n_frames']} frames)  —  cup-near-mouth dominates and is least confounded",
                 fontsize=12.5, weight="bold", y=1.02)
    fig.tight_layout()
    return save(fig, "fig9_failure_factors.png")


def fig_near_mouth_all():
    """All-669 confirmation: failure rate vs cup displacement-from-rest (mocap proxy
    for cup-near-mouth). Near-perfect separation: ~0% at rest -> 31% at the mouth."""
    pa = json.load(open(CACHE / "failure_factors_proxy_all.json"))
    cur = pa["disp_fail_curve"]
    xc = [f"{r['lo']:.0f}–{r['hi']:.0f}" for r in cur]
    yc = [r["fail"] for r in cur]
    fig, ax = plt.subplots(figsize=(9.5, 4.6))
    bars = ax.bar(range(len(xc)), yc, color=C_VIDEO, edgecolor="#333")
    for b, r in zip(bars, cur):
        ax.text(b.get_x() + b.get_width() / 2, r["fail"] + 0.006,
                f"{r['fail']:.1%}\nn={r['n']//1000}k", ha="center", fontsize=8.5)
    ax.set_xticks(range(len(xc))); ax.set_xticklabels(xc, fontsize=9)
    ax.set_xlabel("cup displacement from rest (mm)  →  lifted toward the mouth")
    ax.set_ylabel("failure rate")
    ax.set_ylim(0, max(yc) * 1.2)
    ax.set_title(f"All {pa['n_reps']} reps: failure is gated by cup-near-mouth\n"
                 f"{pa['share_failures_near_mouth']:.0%} of failures occur with the cup lifted "
                 f"(>200 mm); ~0% at rest", weight="bold", fontsize=12)
    return save(fig, "fig10_near_mouth_all.png")


def fig_speed_dist_interaction():
    """Speed x distance interaction: slow worse NEAR the mouth, fast worse FAR.
    Shown for RTS and raw consensus — the sign-flip is in BOTH (real, not a KF
    artifact), while the overall speed level differs (the KF lag part)."""
    si = json.load(open(CACHE / "speed_dist_interaction.json"))
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.4), sharey=True)
    for ax, est, title in zip(axes, ("consensus", "rts"),
                              ("RAW consensus (no filter)", "RTS (smoothed)")):
        rows = si[est]
        x = np.arange(len(rows)); w = 0.38
        slow = [r["slow"] for r in rows]; fast = [r["fast"] for r in rows]
        ax.bar(x - w / 2, slow, w, color="#4c78c8", edgecolor="#333", label="slow")
        ax.bar(x + w / 2, fast, w, color="#e45756", edgecolor="#333", label="fast")
        ax.set_xticks(x); ax.set_xticklabels([r["band"] for r in rows], fontsize=9)
        ax.set_xlabel("cup displacement from rest (mm)\nfar ← → near mouth")
        ax.set_title(title, fontsize=11, weight="bold")
        ax.axhline(0, color="#333", lw=0.8)
    axes[0].set_ylabel("failure rate")
    axes[0].legend(title="cup speed")
    fig.suptitle("Speed × position: slow fails more AT the mouth, fast fails more in transport "
                 "— present in the RAW track (real, not a KF artifact)",
                 fontsize=12, weight="bold", y=1.02)
    fig.tight_layout()
    return save(fig, "fig11_speed_dist.png")


def fig_interp_models():
    """Drinking-phase error vs mocap GT, all on the same 669 reps. KF tuning and
    occlusion-aware hold plateau at the ~20mm apex floor. Two genuinely-different
    methods were tried and BOTH failed on matched reps: a GP (broke time-sync) and
    a learned drink-shape prior (replaced good short-gap KF interpolation with a
    generic template → ~40% WORSE on the reps it touched). The apex is not a fillable
    gap: occlusion is near-universal but mostly confident-WRONG consensus (≥2 cams
    agree on the hand/glass), which no interpolator or prior can fix."""
    ti = json.load(open(CACHE / "tune_interp.json"))
    labels = ["baseline\nKF", "tuned\nKF", "occ-aware\nKF"]
    overall = [ti["baseline"]["median_rms"], ti["tuned_best"]["median_rms"],
               ti["occ_aware"]["median_rms"]]
    drink = [ti["baseline"]["drink_median"], ti["tuned_best"]["drink_median"],
             ti["occ_aware"]["drink_median"]]
    x = np.arange(len(labels)); w = 0.36
    fig, ax = plt.subplots(figsize=(9, 4.8))
    ax.bar(x - w / 2, overall, w, color="#2a9d8f", edgecolor="#333", label="overall")
    ax.bar(x + w / 2, drink, w, color=C_VIDEO, edgecolor="#333", label="drinking phase")
    for xi, (o, d) in enumerate(zip(overall, drink)):
        ax.text(xi - w / 2, o + 0.3, f"{o:.1f}", ha="center", fontsize=9)
        ax.text(xi + w / 2, d + 0.3, f"{d:.0f}", ha="center", fontsize=9, weight="bold")
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=11)
    ax.set_ylabel("median error vs mocap (mm)  —  all 669 reps")
    ax.set_ylim(0, max(drink) * 1.2)
    ax.set_title("Nothing cracks the ~20 mm apex floor — the failure is confident-WRONG\n"
                 "occlusion, not a fillable gap (interpolators & a learned prior both fail)",
                 weight="bold", fontsize=11.5)
    ax.legend(loc="center right")
    ax.text(0.5, -0.25, "(GP and a learned drink-shape prior both tried on matched reps — "
            "neither helped; the shape prior was ~40% worse where it fired)",
            transform=ax.transAxes, ha="center", fontsize=8, style="italic", color="#888")
    return save(fig, "fig12_interp_models.png")


figs = [fig_pipeline(), fig_method_example(), fig_rms_hist(),
        fig_class_bar(), fig_scatter(), fig_per_participant(),
        fig_fracfail_by_phase(), fig_within_phase(), fig_failure_factors(),
        fig_near_mouth_all(), fig_speed_dist_interaction(), fig_interp_models()]


# ----------------------------------------------------------------------- pptx
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide(title, img, subtitle=""):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        p.font.size = Pt(26); p.font.bold = True
        p.font.color.rgb = RGBColor(0x26, 0x46, 0x53)
        if subtitle:
            sp = tf.add_paragraph(); sp.text = subtitle
            sp.font.size = Pt(13); sp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        from PIL import Image
        iw, ih = Image.open(img).size
        maxw, maxh = Inches(12.3), Inches(5.7)
        scale = min(maxw / iw, maxh / ih)
        w, h = int(iw * scale), int(ih * scale)
        s.shapes.add_picture(str(img), int((prs.slide_width - w) / 2), Inches(1.5),
                             width=w, height=h)
        return s

    def bullets(title, items, subtitle=""):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(1.0))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        p.font.size = Pt(28); p.font.bold = True
        p.font.color.rgb = RGBColor(0x26, 0x46, 0x53)
        if subtitle:
            sp = tf.add_paragraph(); sp.text = subtitle
            sp.font.size = Pt(14); sp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        body = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.3))
        bf = body.text_frame; bf.word_wrap = True
        for i, it in enumerate(items):
            lvl = 0
            if isinstance(it, tuple):
                it, lvl = it
            par = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            par.text = ("– " if lvl else "• ") + it
            par.level = lvl
            par.font.size = Pt(18 if lvl == 0 else 15)
            par.font.color.rgb = RGBColor(0x22, 0x22, 0x22) if lvl == 0 else RGBColor(0x55, 0x55, 0x55)
            par.space_after = Pt(8 if lvl == 0 else 4)
        return s

    # title
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Validating video cup-tracking against QTM mocap"
    p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = RGBColor(0x26, 0x46, 0x53)
    sp = tf.add_paragraph()
    sp.text = f"{N} drinking reps · 22 participants · median {np.median(ir):.1f} mm agreement"
    sp.font.size = Pt(20); sp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ---- methodology section (detailed) ----
    bullets("Methodology — data & ground truth",
            [("Two independent recordings of the same drinking reps:", 0),
             ("QTM optical mocap — 4 markers on the cup, 100/120 Hz, sub-mm 6-DoF "
              "pose. This is the ground truth.", 1),
             ("Video cup track — 10 calibrated cameras, per-camera detections fused "
              "by ≥3-cam gated triangulation (consensus), then a Kalman filter + RTS "
              "smoother → one 3D cup centroid in mm.", 1),
             ("Phases are segmented from the MOCAP trajectory (van Andel protocol: "
              "6 Hz filter, velocity gates) → rest / forward-transport / drinking / "
              "back-transport / rest. GT-defined, so independent of video error.", 0),
             ("669 valid drinking reps across 22 participants.", 0)],
            "What we compare and what counts as truth")

    bullets("Methodology — pairing & temporal sync",
            [("Each video rep must be matched to its mocap take, then time-aligned:", 0),
             ("Pairing: mocap recorded extra/aborted takes, and drinking speed-curves "
              "are near-identical, so duration matching mis-registers. We pair by the "
              "order-preserving assignment that minimises 3D Kabsch fit residual "
              "(2 mm vs 30 mm is unambiguous).", 1),
             ("Sync: cross-correlate the two cup speed-curves; the lag at peak "
              "correlation aligns them in time. Reps below 0.7 correlation are rejected.", 1),
             ("Ground-truth quality gate: mocap takes that are dead/static or have "
              "non-physical jumps are dropped before scoring, so a bad GT can never be "
              "charged as tracking error.", 0)],
            "How a video rep is put in correspondence with mocap")

    bullets("Methodology — error metric & failure",
            [("After sync, fit a robust rigid transform (Huber-weighted Kabsch) so we "
              "measure trajectory SHAPE, not absolute placement.", 0),
             ("Per frame: distance between the transformed video track and the mocap "
              "cup (mm). Two summary numbers per rep:", 0),
             ("inlier RMS — fidelity on the frames that agree (tracking accuracy).", 1),
             ("frac_fail — fraction of the trajectory >50 mm off (coverage of failure).", 1),
             ("A frame 'fails' when the mocap–video distance exceeds 50 mm. Each "
              "failing frame is attributed to its mocap phase to localise WHERE "
              "tracking breaks.", 0),
             ("Driver analysis: logistic regression of per-frame failure on physical "
              "factors (cup-near-mouth, hand-on-cup, speed), with cluster-bootstrap "
              "CIs (resample reps, not frames).", 0)],
            "From aligned trajectories to a failure label")

    slide("Methodology — pipeline at a glance", figs[0],
          "Mocap (sub-mm) is the ground truth; we measure how close the video track lands.")
    slide("How each rep is scored", figs[1],
          "Two numbers: fidelity where it tracks, and how much of the trajectory fails.")
    slide("Result: agreement to a few millimetres", figs[2],
          f"Median {np.median(ir):.1f} mm, p90 {np.percentile(ir,90):.1f} mm over {N} reps.")
    slide("Result: accuracy classes", figs[3],
          "72% clean, 28% localized (brief occluded apex), 0 broken.")
    slide("Result: fidelity vs failure coverage", figs[4],
          "Every rep tracks <15 mm; the class only reflects the occluded drink apex.")
    slide("Result: consistent across participants", figs[5])
    slide("Result: failure is confined to the drinking phase", figs[6],
          "Rest phases ~0% fail (2–3 mm); only the occluded mouth dwell drives error.")
    slide("Result: where within the phase failure peaks", figs[7],
          "Failure rate peaks mid-drinking — the apex where the cup reaches the mouth.")
    slide("Result: what drives failure (untangled)", figs[8],
          "Cup-near-mouth dominates; hand-on-cup is a gate (fail only when on); speed is non-monotonic (slow@mouth, fast@transport).")
    slide("Result: confirmed on all 669 reps", figs[9],
          "Mocap proxy (cup displacement-from-rest, validated r≈−0.99 vs real cup→mouth): failure rises ~0%→31% with reach to the mouth.")
    slide("Result: speed × position (and the KF's role)", figs[10],
          "Slow fails more at the mouth, fast in transport — the sign-flip survives in the RAW track, so the interaction is real; the overall speed level is partly KF lag.")
    slide("Can we improve the interpolator?", figs[11],
          "No — KF tuning, occlusion-hold, a GP, and a learned drink-shape prior all fail to beat ~20 mm (the shape prior was worse on matched reps). Occlusion at the mouth is near-universal and mostly confident-WRONG consensus, not a fillable gap; the fix is better input (camera view / appearance gating), not a better filter or prior.")

    # footnote slide
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Footnote: rep–take pairing"
    p.font.size = Pt(26); p.font.bold = True
    p.font.color.rgb = RGBColor(0x26, 0x46, 0x53)
    for line in [
        "Mocap recorded more takes than the video annotation logged (extra / aborted captures).",
        "Drinking speed curves are near-identical, so duration-based pairing mis-registered one side (P02 right).",
        "Fix: pair each rep to its take by the order-preserving assignment that minimises Kabsch fit residual (2 mm vs 30 mm is unambiguous).",
        "Effect: P02 broken reps 2 → 0; cohort broken 2 → 0; all other participants unchanged.",
    ]:
        b = tf.add_paragraph(); b.text = "• " + line
        b.font.size = Pt(16); b.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        b.space_after = Pt(10)

    out = OUT / "qtm_accuracy.pptx"
    prs.save(out)
    print("wrote", out)


build_pptx()
print("\nDone. Figures + deck in", OUT)
