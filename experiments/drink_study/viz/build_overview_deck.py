"""Concise WHOLE-PIPELINE deck (lab/supervisor update) -- ~10 slides, one claim each.
Reuses the existing slides/*.png figures + the new fig13_dispute_strip. Writes a NEW
file (slides/pipeline_overview.pptx) so it never clobbers the open qtm_accuracy.pptx.

    python experiments/drink_study/build_overview_deck.py
"""
from __future__ import annotations
import json
from pathlib import Path
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

HERE = Path(__file__).resolve().parent
SL = HERE / "slides"
from _paths import CACHE

# headline number from the validation cache
d = json.load(open(CACHE / "qtm_align.json"))
reps = [r for p in d.values() if isinstance(p, dict) and p.get("ok") for r in p["reps"]]
ir = np.array([r["inlier_rms_mm"] for r in reps])
N = len(reps)
MED = np.median(ir)

NAVY = RGBColor(0x26, 0x46, 0x53)
GREY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def _title(s, title, subtitle=""):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = NAVY
    if subtitle:
        sp = tf.add_paragraph(); sp.text = subtitle
        sp.font.size = Pt(13); sp.font.color.rgb = GREY


def fig_slide(title, img, subtitle="", top=1.55, maxh=5.6):
    s = prs.slides.add_slide(blank)
    _title(s, title, subtitle)
    img = SL / img
    iw, ih = Image.open(img).size
    maxw = Inches(12.3); maxhi = Inches(maxh)
    scale = min(maxw / iw, maxhi / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(img), int((prs.slide_width - w) / 2), Inches(top), width=w, height=h)
    return s


def bullets_slide(title, items, subtitle=""):
    s = prs.slides.add_slide(blank)
    _title(s, title, subtitle)
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.3))
    bf = body.text_frame; bf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        par = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        par.text = ("– " if lvl else "• ") + it
        par.level = lvl
        par.font.size = Pt(19 if lvl == 0 else 15)
        par.font.color.rgb = DARK if lvl == 0 else GREY
        par.space_after = Pt(9 if lvl == 0 else 4)
    return s


# ---- 1. Title ----
s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.3), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Markerless cup tracking, end to end"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
sp = tf.add_paragraph()
sp.text = "Detection → 3D fusion → tracking accuracy → phase segmentation"
sp.font.size = Pt(20); sp.font.color.rgb = GREY
sp2 = tf.add_paragraph()
sp2.text = f"Validated on {N} drinking reps · 22 participants · median {MED:.1f} mm vs QTM mocap"
sp2.font.size = Pt(16); sp2.font.color.rgb = GREY

# ---- 2. Pipeline ----
fig_slide("The pipeline", "fig1_pipeline.png",
          "10 cameras → per-camera detection → ≥3-cam consensus + KF/RTS → one 3D cup track, "
          "validated frame-by-frame against sub-mm mocap.")

# ---- 3. Accuracy ----
fig_slide("It tracks to a few millimetres", "fig3_rms_hist.png",
          f"Per-rep inlier RMS vs mocap: median {MED:.1f} mm, p90 {np.percentile(ir,90):.1f} mm "
          f"over {N} reps — and consistent across all 22 participants.")

# ---- 4. ...except the apex ----
fig_slide("…except at the drinking apex", "fig8_within_phase.png",
          "Failure is one continuous arc: ~0% in rest/transport, peaking mid-drink where the "
          "cup reaches the mouth. Everything downstream is about that apex.")

# ---- 5. Why it fails ----
fig_slide("Why it fails: cup at the mouth", "fig9_failure_factors.png",
          "Cup-near-mouth dominates (and is least confounded); hand-on-cup is a gate; speed is "
          "non-monotonic. At the mouth the cameras are occluded and the consensus is confidently WRONG.")

# ---- 6. Apex not fillable ----
fig_slide("The apex is not a fillable gap", "fig12_interp_models.png",
          "KF-tuning, occlusion-hold, a GP and a learned shape-prior all stall at the ~20 mm floor — "
          "the failure is confident-wrong occlusion, not missing data.")

# ---- 7. ...but the gap is (velocity-fill) ----
fig_slide("But short gaps ARE fillable — predict movement", "fig_tcn_gapfill.png",
          "Velocity-fill TCN predicts the cup's MOVEMENT (not position) across occlusion: "
          "−13% median / −33% p90 at the apex vs KF coast. (Only LOPO-surviving gains kept.)")

# ---- 8. Segmentation: gate → learned ----
fig_slide("Phase segmentation: from a gate to a learned model", "fig14_segmenter.png",
          "Cup-only dwell segmentation. Learning beats the gate on the TYPICAL rep once "
          "(median 133→~110), then flat; richer 3D-direction + occlusion features don't move "
          "the median further — they cut the EXTREME tail where a scalar gate structurally "
          "fails (p99 1056→700, max 3617→2067). Honest trade: hybrid over-extends ~26 easy "
          "reps — ship is a judgment call, gate still in production.",
          top=1.85, maxh=4.9)

# ---- 9. The truth is imperfect ----
fig_slide("The catch: our 'truth' is itself imperfect", "fig13_dispute_strip.png",
          "P23 cam4. Green = frames the HYBRID calls drinking but the mocap speed-gate 'truth' "
          "excludes. The cup is at the mouth in ALL of them — the model is MORE right than its label.",
          top=2.3, maxh=3.6)
# add the interpretation line lower
s9 = prs.slides[-1]
tb = s9.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(12), Inches(1.1))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Truth = cup-nearly-STILL (speed gate, no mouth marker); real drinking = cup-AT-MOUTH, "
          "which brackets the still-phase. Some 'regressions' are the model catching onset/offset "
          "the label misses. True validation needs a mouth/face marker.")
p.font.size = Pt(13); p.font.color.rgb = GREY

# ---- 10. Where it stands / open ----
bullets_slide("Where it stands & what's open",
    [("Tracking: markerless video matches sub-mm mocap to a few mm; 0 broken reps.", 0),
     ("The only failure is the occluded drinking apex — confident-wrong consensus, not a "
      "fillable gap; better INPUT (view/appearance gating) is the lever, not a better filter.", 0),
     ("Gap-fill: velocity-fill TCN recovers short apex gaps by predicting movement.", 0),
     ("Segmentation: hybrid learned segmenter beats the tuned gate on mean + tail (a trade).", 0),
     ("Open — the real frontier:", 0),
     ("Add a mouth/face marker (van Andel 15%-of-steady-state) for a behaviourally correct "
      "dwell truth — our current speed-gate label is provably late.", 1),
     ("Decide ship: hybrid (worst-case/mean) vs tuned gate (never breaks easy reps).", 1),
     ("Port consensus-anchored KF into run_pipeline.", 1)],
    "Summary and next steps")

out = SL / "pipeline_overview.pptx"
prs.save(out)
print("wrote", out, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
